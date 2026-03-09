import os
import json
import google.generativeai as genai
from docx import Document as DocxDocument
from pptx import Presentation
import base64

genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel("gemini-1.5-flash")

# Static section templates per artifact type
SECTION_TEMPLATES = {
    "SOW": [
        ("SOW_ProjectOverview", "Project Overview"),
        ("SOW_Objectives", "Objectives"),
        ("SOW_InScope", "In Scope"),
        ("SOW_OutOfScope", "Out of Scope"),
        ("SOW_Assumptions", "Assumptions"),
        ("SOW_Dependencies", "Dependencies"),
        ("SOW_Deliverables", "Deliverables"),
        ("SOW_Milestones", "Milestones"),
        ("SOW_Timeline", "Timeline"),
        ("SOW_RolesResponsibilities", "Roles & Responsibilities"),
        ("SOW_Pricing", "Pricing"),
        ("SOW_PaymentTerms", "Payment Terms"),
        ("SOW_ChangeControl", "Change Control"),
        ("SOW_RisksMitigations", "Risks & Mitigations"),
        ("SOW_SLAs", "SLAs"),
        ("SOW_AcceptanceCriteria", "Acceptance Criteria"),
    ],
    "Proposal": [
        ("Proposal_AccountOverview", "Account Overview"),
        ("Proposal_CurrentState", "Current State"),
        ("Proposal_ProposedSolution", "Proposed Solution"),
        ("Proposal_ScopeSummary", "Scope Summary"),
        ("Proposal_BusinessBenefits", "Business Benefits"),
        ("Proposal_ImplementationApproach", "Implementation Approach"),
        ("Proposal_Timeline", "Timeline"),
        ("Proposal_TeamProfiles", "Team Profiles"),
        ("Proposal_PricingSummary", "Pricing Summary"),
        ("Proposal_CommercialTerms", "Commercial Terms"),
        ("Proposal_CaseStudies", "Case Studies"),
    ],
    "DesignDocument": [
        ("Design_BusinessCapabilities", "Business Capabilities"),
        ("Design_UserStories", "User Stories"),
        ("Design_DataModel", "Data Model"),
        ("Design_IntegrationDesign", "Integration Design"),
        ("Design_ProcessDesign", "Process Design"),
        ("Design_SecurityModel", "Security Model"),
        ("Design_NFRs", "Non-Functional Requirements"),
        ("Design_EnvironmentDeployment", "Environment & Deployment"),
        ("Design_ConfigVsCustom", "Config vs Custom"),
        ("Design_OpenQuestions", "Open Questions"),
        ("Design_Risks", "Risks"),
    ],
    "ProcessMap": [
        ("ProcessMap_Journey", "Journey"),
        ("ProcessMap_Lanes", "Lanes"),
    ],
    "DiscoveryNotes": [
        ("Discovery_BusinessContext", "Business Context"),
        ("Discovery_PainPoints", "Pain Points"),
        ("Discovery_UseCases", "Use Cases"),
        ("Discovery_SystemsLandscape", "Systems Landscape"),
        ("Discovery_DataNeeds", "Data Needs"),
        ("Discovery_RisksConstraints", "Risks & Constraints"),
    ],
    "Other": [
        ("Other_Section", "General Section"),
    ],
}

ARTIFACT_SUB_TYPES = {
    "SOW": "SOW_Document",
    "Proposal": "Proposal_Document",
    "DesignDocument": "Design_Document",
    "ProcessMap": "ProcessMap_Document",
    "DiscoveryNotes": "Discovery_Document",
    "Other": "Other_Document",
}


def extract_text_from_docx(path):
    doc = DocxDocument(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def extract_text_from_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)


def extract_sections_with_llm(text, artifact_type, filename, is_image=False, image_path=None):
    sections = SECTION_TEMPLATES.get(artifact_type, SECTION_TEMPLATES["Other"])
    section_names = [s[1] for s in sections]
    section_subtypes = [s[0] for s in sections]

    prompt = f"""You are a knowledge extraction assistant.
You will be given a document of type: {artifact_type}.
Extract content for each of these sections: {', '.join(section_names)}.
Return a JSON object where keys are the section names and values are the extracted content (2-5 sentences max per section).
If a section is not present in the document, return an empty string for that section.
Return ONLY valid JSON, no other text, no markdown fences.

Document: {filename}

Content:
{text[:6000]}"""

    try:
        if is_image and image_path:
            import PIL.Image
            img = PIL.Image.open(image_path)
            response = model.generate_content([prompt, img])
        else:
            response = model.generate_content(prompt)

        raw = response.text.strip()
        # Clean markdown fences if present
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        raw = raw.strip()

        extracted = json.loads(raw)
    except Exception as e:
        print(f"LLM extraction error for {filename}: {e}")
        extracted = {}

    result = []
    for subtype, name in zip(section_subtypes, section_names):
        content = extracted.get(name, "")
        result.append((subtype, name, content))

    return result


def ingest_source(source, db):
    artifact_type = source.artifact_type
    path = source.stored_path
    filename = source.filename

    is_image = False
    text = ""

    ext = filename.lower().split(".")[-1]
    if ext == "docx":
        text = extract_text_from_docx(path)
    elif ext == "pptx":
        text = extract_text_from_pptx(path)
    elif ext in ("png", "jpg", "jpeg"):
        is_image = True
    else:
        text = ""

    sections = extract_sections_with_llm(
        text, artifact_type, filename,
        is_image=is_image,
        image_path=path if is_image else None
    )

    return sections